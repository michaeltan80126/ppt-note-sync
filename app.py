"""
PPT Note Sync - 讲稿同步助手
使用 Streamlit 构建的Web应用
将逐字稿自动匹配到PPT的演讲者备注中
"""

import streamlit as st
import re
import io
from datetime import datetime
from pptx import Presentation
from docx import Document

# 页面配置
st.set_page_config(
    page_title="PPT Note Sync - 讲稿同步助手",
    page_icon="📊",
    layout="centered"
)

# 自定义CSS
st.markdown("""
<style>
    .main {
        background-color: #f5f5f5;
    }
    .stButton>button {
        width: 100%;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        padding: 12px;
        border-radius: 8px;
        font-size: 16px;
        font-weight: 500;
    }
    .stButton>button:hover {
        transform: translateY(-2px);
        box-shadow: 0 5px 20px rgba(102, 126, 234, 0.4);
    }
    .info-box {
        background: #fff3e0;
        border-left: 4px solid #ff9800;
        padding: 15px;
        margin-bottom: 20px;
        border-radius: 0 8px 8px 0;
    }
    .success-box {
        background: #e8f5e9;
        border-left: 4px solid #4caf50;
        padding: 15px;
        border-radius: 0 8px 8px 0;
    }
    .header {
        text-align: center;
        padding: 20px;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 10px;
        margin-bottom: 30px;
    }
</style>
""", unsafe_allow_html=True)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """从DOCX文件中提取纯文本"""
    try:
        doc = Document(io.BytesIO(file_bytes))
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        st.error(f"无法读取DOCX文件: {str(e)}")
        return ""


def extract_text_from_txt(file_bytes: bytes) -> str:
    """从TXT文件中提取文本"""
    try:
        # 尝试多种编码
        for encoding in ['utf-8', 'gbk', 'gb2312', 'utf-16']:
            try:
                return file_bytes.decode(encoding)
            except:
                continue
        st.error("无法识别文件编码")
        return ""
    except Exception as e:
        st.error(f"无法读取TXT文件: {str(e)}")
        return ""


def parse_script(script_content: str) -> dict:
    """
    解析逐字稿内容，返回 {slide_index: content} 的字典
    使用 ### Slide X 作为分隔符
    """
    # 移除HTML标签
    script_content = re.sub(r'<[^>]+>', '', script_content)
    
    # 使用正则表达式匹配 ### Slide X
    pattern = r'###\s*Slide\s*(\d+)'
    matches = list(re.finditer(pattern, script_content, re.IGNORECASE))
    
    result = {}
    
    if not matches:
        # 如果没有找到分隔符，将整个内容作为第一页
        result[1] = script_content.strip()
        return result
    
    # 遍历每个匹配
    for i, match in enumerate(matches):
        slide_num = int(match.group(1))
        start = match.end()
        
        # 获取下一个分隔符的位置
        if i + 1 < len(matches):
            end = matches[i + 1].start()
        else:
            end = len(script_content)
        
        # 提取内容
        content = script_content[start:end].strip()
        result[slide_num] = content
    
    return result


def process_ppt(pptx_file, script_data: dict) -> bytes:
    """
    处理PPT文件，将逐字稿内容注入到演讲者备注
    返回处理后的PPTX文件字节
    """
    try:
        # 读取上传的PPT文件
        prs = Presentation(pptx_file)
        
        # 遍历每一页幻灯片
        processed_count = 0
        for slide_num, content in script_data.items():
            # PPT的slide索引从0开始
            slide_index = slide_num - 1
            
            if slide_index < 0 or slide_index >= len(prs.slides):
                continue
            
            slide = prs.slides[slide_index]
            
            # 获取或创建备注
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            
            # 清空原有备注并写入新内容
            text_frame.clear()
            text_frame.text = content
            processed_count += 1
        
        # 保存到内存
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.getvalue(), processed_count
        
    except Exception as e:
        st.error(f"处理PPT文件时出错: {str(e)}")
        return None, 0


# 主界面
st.markdown("""
<div class="header">
    <h1>📊 PPT Note Sync</h1>
    <p>将逐字稿自动匹配到PPT演讲者备注</p>
</div>
""", unsafe_allow_html=True)

# 说明框
st.markdown("""
<div class="info-box">
    <h4>📝 逐字稿格式说明</h4>
    <p>请使用 <code>### Slide X</code> 作为每一页的分隔标记。例如：</p>
    <p style="margin-top: 8px;">
        <code>### Slide 1</code><br>
        第一页的内容...<br>
        <code>&lt;break time="1.0s" /&gt;</code><br>
        继续内容...<br><br>
        <code>### Slide 2</code><br>
        第二页的内容...
    </p>
</div>
""", unsafe_allow_html=True)

# 文件上传
col1, col2 = st.columns(2)

with col1:
    st.subheader("📄 上传 PPT 文件")
    pptx_file = st.file_uploader("选择PPT文件", type=['pptx'], key='pptx')

with col2:
    st.subheader("📝 上传逐字稿")
    script_file = st.file_uploader("选择逐字稿文件", type=['txt', 'docx'], key='script')

# 处理按钮
if st.button("🚀 开始处理"):
    if not pptx_file:
        st.error("请上传PPT文件")
    elif not script_file:
        st.error("请上传逐字稿文件")
    else:
        with st.spinner('正在处理文件，请稍候...'):
            try:
                # 读取逐字稿
                script_content = script_file.read()
                
                if script_file.name.endswith('.docx'):
                    script_text = extract_text_from_docx(script_content)
                else:
                    script_text = extract_text_from_txt(script_content)
                
                if not script_text:
                    st.error("无法读取逐字稿内容")
                else:
                    # 解析逐字稿
                    script_data = parse_script(script_text)
                    
                    # 处理PPT
                    result, count = process_ppt(pptx_file, script_data)
                    
                    if result:
                        # 生成输出文件名
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        output_filename = f"PPT备注同步_{timestamp}.pptx"
                        
                        # 显示成功信息
                        st.markdown("""
                        <div class="success-box">
                            <h4>✅ 处理完成！</h4>
                            <p>已成功将逐字稿内容添加到PPT的演讲者备注中</p>
                            <p>共处理 <b>{}</b> 页幻灯片</p>
                        </div>
                        """.format(count), unsafe_allow_html=True)
                        
                        # 下载按钮
                        st.download_button(
                            label="📥 下载处理后的PPT",
                            data=result,
                            file_name=output_filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
            except Exception as e:
                st.error(f"处理失败: {str(e)}")

# 页脚
st.markdown("---")
st.markdown("<p style='text-align: center; color: #999;'>PPT Note Sync - 讲稿同步助手</p>", unsafe_allow_html=True)
